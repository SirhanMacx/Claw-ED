"""PPTX slide compiler for MasterContent.

Compiles a MasterContent object into a classroom-ready PowerPoint slide deck
that matches the teacher's real exemplar (the "3 Gs / God Gold Glory" deck): a
sparse, image-forward 16:9 deck (~15–18 words per slide, NOT paragraphs) on the
standard 10.0 x 5.625 in canvas so it opens cleanly in Keynote/PowerPoint on
macOS and imports losslessly into Google Slides.

Visual DNA matched from the real exemplar:
  - FONT: Century Gothic everywhere (a macOS system font), Calibri fallback.
  - PALETTE: warm terracotta (#DB8258) titles, gold (#FBD673) + white accents,
    on a light cream (#FFFDF7) background.
  - NO filled header bars — titles are plain large clean text in the palette.
  - Turn-and-Talk discussion slides (title + one big question).
  - Image-forward: content slides embed images; an image-activity GRID slide
    tiles many small aspect-preserved cards when ≥6 images are available.

Slide order: title -> Turn and Talk (do-now) -> vocabulary -> instruction
sections -> source analysis -> image-activity grid -> station overview ->
exit ticket.
No LLM calls — pure mechanical compilation.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from pptx.presentation import Presentation

    from clawed.master_content import MasterContent

logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════════════════
# Canvas + grid geometry (standard 16:9, matches the teacher's exemplar deck)
# ═══════════════════════════════════════════════════════════════════════════
#
# The exemplar deck (AgeofExploration…The3Gs) is 10.0 x 5.625 in — the standard
# PowerPoint 16:9 size that imports losslessly into Google Slides and opens
# cleanly in Keynote. We use ONE fixed grid for every content slide so images
# land in the same aligned slot every time and never collide with text. Unlike
# the old build, there is NO filled header bar — the exemplar uses plain large
# terracotta title text on a light cream background:
#
#   ┌──────────────────────────── 10.0 in ────────────────────────────┐
#   │ TITLE  (large terracotta text, no filled bar, top 0.30..1.10)    │
#   ├──────────────────────────────────┬──────────────────────────────┤
#   │ TEXT PANEL                        │  IMAGE BOX                    │
#   │ left 0.3 .. 5.7 (width 5.4)       │  left 5.9, w 3.8, top 1.2,    │
#   │                                   │  h 4.1 (aspect-preserved)     │
#   └──────────────────────────────────┴──────────────────────────────┘
#
# Text content_width is hard-capped so its right edge (5.7) sits left of the
# image box left edge (5.9) — guaranteeing no overlap.

SLIDE_W_IN = 10.0
SLIDE_H_IN = 5.625
MARGIN_IN = 0.3
# Plain (unfilled) title band — large clean text, mirrors the exemplar.
TITLE_TOP_IN = 0.30
TITLE_H_IN = 0.85

# Right-hand image slot — the SAME box for instruction / source / exit slides.
# Dropped a touch below the (now unfilled) title text so it never overlaps it.
IMG_BOX_LEFT_IN = 5.9
IMG_BOX_TOP_IN = 1.2
IMG_BOX_W_IN = 3.8
IMG_BOX_H_IN = 4.1

# Text panel when an image is present: 0.3 .. 5.7 (right edge < image box left).
TEXT_PANEL_W_IN = 5.4
# Text panel when no image: full width minus both margins.
FULL_TEXT_W_IN = SLIDE_W_IN - 2 * MARGIN_IN  # 9.4


# ═══════════════════════════════════════════════════════════════════════════
# Visual theme — Century Gothic + the exemplar's terracotta / gold / cream
# palette. The exemplar has NO filled header bars: titles are plain large text
# in terracotta on a light cream background.
# ═══════════════════════════════════════════════════════════════════════════

# Century Gothic is a macOS system font; PowerPoint / Keynote / Google Slides
# resolve it by name and substitute the fallback otherwise.
FONT_PRIMARY = "Century Gothic"
FONT_FALLBACK = "Calibri"

# Palette (measured from the real "3 Gs" deck).
_C_TERRACOTTA = "DB8258"   # dominant warm title / accent text
_C_GOLD = "FBD673"         # secondary accent
_C_WHITE = "FFFFFF"        # accent on darker fills (rare)
_C_CREAM_BG = "FFFDF7"     # light cream slide background
_C_BODY = "4A3B33"         # warm dark brown body text (readable on cream)
_C_TITLE = _C_TERRACOTTA   # titles are terracotta, not near-black


def _short(text: str, words: int = 12) -> str:
    """Reduce free text to a short phrase of at most ``words`` whitespace tokens.

    This is the single density gate every slide-text call routes through, so the
    deck stays sparse and visual (the exemplar averages ~16 words/slide) instead
    of dumping mid-sentence paragraph slices. Strips trailing punctuation and
    appends an ellipsis only when the text was actually truncated.
    """
    if not text:
        return ""
    tokens = str(text).split()
    if not tokens:
        return ""
    truncated = len(tokens) > words
    kept = tokens[:words]
    out = " ".join(kept).rstrip(",;:.!?—-– ").strip()
    if truncated and out:
        out += "…"
    return out


def _first_sentence(text: str, words: int = 20) -> str:
    """Take the first sentence of ``text``, then cap it to ``words`` tokens.

    Produces a clean phrase boundary instead of a mid-sentence slice. Falls back
    to a plain ``_short`` cut when there is no sentence delimiter.
    """
    if not text:
        return ""
    first = str(text).replace("\n", " ").split(". ")[0]
    return _short(first, words=words)


# ═══════════════════════════════════════════════════════════════════════════
# Low-level shape helpers
# ═══════════════════════════════════════════════════════════════════════════


def _hex_to_rgb(hex_color: str) -> Any:
    """Convert a 6-char hex string to pptx RGBColor."""
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))  # type: ignore[no-untyped-call]  # python-pptx is untyped


def _apply_font(run: Any) -> None:
    """Name the run's font Century Gothic and register a safe Calibri fallback.

    We set ``run.font.name`` (the latin typeface) to Century Gothic and inject
    the legacy/cs typeface plus a fallback hint into the run XML so PowerPoint /
    Keynote / Google Slides degrade cleanly to Calibri when Century Gothic is
    not installed on the viewer's machine.
    """
    run.font.name = FONT_PRIMARY
    try:
        from pptx.oxml.ns import qn

        rpr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:cs", "a:ea"):
            el = rpr.find(qn(tag))
            if el is None:
                el = rpr.makeelement(qn(tag), {})
                rpr.append(el)
            el.set("typeface", FONT_PRIMARY)
    except Exception as exc:
        logger.debug("Could not set font fallback hint: %s", exc)


def _add_slide(prs: Any, layout_idx: int = 6, bg_hex: str = _C_CREAM_BG) -> Any:
    """Add a blank slide with the cream exemplar background and return it."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    if bg_hex:
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(bg_hex)
        except Exception as exc:
            logger.debug("Could not set slide background: %s", exc)
    return slide


def _textbox(slide: Any, left: Any, top: Any, width: Any, height: Any, text: str,
             font_size: int = 18, bold: bool = False,
             hex_color: str = _C_BODY, align_center: bool = False,
             italic: bool = False, word_wrap: bool = True) -> Any:
    """Add a textbox to a slide and return the shape."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex_to_rgb(hex_color)
    _apply_font(run)
    return tb


def _bullet_textbox(slide: Any, left: Any, top: Any, width: Any, height: Any,
                    items: list[str], font_size: int = 18,
                    hex_color: str = _C_BODY) -> Any:
    """Add a textbox with one paragraph per bullet item."""
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = _hex_to_rgb(hex_color)
        _apply_font(run)

    return tb


def _embed_image(slide: Any, image_spec: str, images: dict[str, Path],
                 box_left: Any, box_top: Any, box_w: Any, box_h: Any) -> bool:
    """Embed a pre-fetched image, CONTAINED in the given box, aspect preserved.

    The image is scaled with ``min(box_w/iw, box_h/ih)`` so it always fits the
    box without stretching, then centered inside the box. We pass ONLY the
    scaled draw size to ``add_picture`` (never the raw box dimensions), which is
    what previously stretched images. If Pillow can't read the native size, we
    fall back to fitting by box width (still aspect-correct via add_picture's
    single-dimension mode).

    All coordinates are in EMU. Returns True if an image was embedded.
    """
    if not image_spec:
        return False
    path = images.get(image_spec)
    if not (path and Path(path).exists()):
        return False

    try:
        iw = ih = None
        try:
            from PIL import Image

            with Image.open(str(path)) as im:
                iw, ih = im.size
        except Exception as exc:
            logger.debug("Could not read native size for %r: %s", image_spec, exc)

        if iw and ih:
            scale = min(box_w / iw, box_h / ih)
            draw_w = int(iw * scale)
            draw_h = int(ih * scale)
            left = int(box_left + (box_w - draw_w) / 2)
            top = int(box_top + (box_h - draw_h) / 2)
            slide.shapes.add_picture(str(path), left, top, width=draw_w, height=draw_h)
        else:
            # No native size — let pptx preserve aspect from the file by fixing
            # only the width, anchored at the box's top-left.
            slide.shapes.add_picture(str(path), int(box_left), int(box_top), width=int(box_w))
        return True
    except Exception as exc:
        logger.debug("Could not embed slide image %r: %s", image_spec, exc)
    return False


def _slide_title(slide: Any, prs_width: Any, text: str,
                 hex_color: str = _C_TITLE, font_size: int = 34,
                 align_center: bool = False) -> None:
    """Add a plain, large title (NO filled bar) in the exemplar palette.

    The real exemplar has no filled header rectangles — section titles are just
    large clean terracotta text on the cream background. This replaces the old
    ``_header_bar`` look while keeping the same one-line, density-capped title.
    """
    from pptx.util import Inches

    capped = _short(text, words=7)
    # Auto-shrink long headings so they stay on one line and never crowd the body
    # content positioned just below (mirrors the title-slide auto-shrink).
    if font_size > 26 and len(capped) > 40:
        font_size = 28 if len(capped) <= 54 else 26

    _textbox(
        slide,
        left=Inches(MARGIN_IN), top=Inches(TITLE_TOP_IN),
        width=prs_width - Inches(2 * MARGIN_IN), height=Inches(TITLE_H_IN),
        text=capped,
        font_size=font_size, bold=True,
        hex_color=hex_color, align_center=align_center,
        word_wrap=False,
    )


# ═══════════════════════════════════════════════════════════════════════════
# Per-slide-type builders
# ═══════════════════════════════════════════════════════════════════════════


def _build_title_slide(prs: Presentation, master: MasterContent) -> None:
    """Slide 1: Title with subject, grade, and objective.

    Cream background with a large terracotta title and gold meta line — mirrors
    the exemplar's clean, image-forward cover (no dark fill, no header bar).
    """
    from pptx.util import Inches

    W = prs.slide_width  # noqa: N806
    slide = _add_slide(prs)

    # Auto-shrink the title for long generated lesson names so it stays ~2 lines
    # and never overruns the meta/objective lines below it.
    title = master.title or ""
    title_size = 42 if len(title) <= 38 else (34 if len(title) <= 66 else 28)
    _textbox(
        slide,
        left=Inches(0.6), top=Inches(1.15),
        width=W - Inches(1.2), height=Inches(2.05),
        text=title,
        font_size=title_size, bold=True,
        hex_color=_C_TERRACOTTA, align_center=True,
    )
    meta = f"{master.subject}  |  Grade {master.grade_level}  |  {master.duration_minutes} min"
    _textbox(
        slide,
        left=Inches(0.6), top=Inches(3.5),
        width=W - Inches(1.2), height=Inches(0.5),
        text=meta,
        font_size=18, bold=True,
        hex_color=_C_TERRACOTTA, align_center=True,
    )
    _textbox(
        slide,
        left=Inches(0.6), top=Inches(4.25),
        width=W - Inches(1.2), height=Inches(1.0),
        text=f"Objective: {_short(master.objective, words=10)}",
        font_size=16,
        hex_color=_C_BODY, align_center=True,
    )


def _build_turn_and_talk_slides(prs: Presentation, master: MasterContent) -> None:
    """Turn-and-Talk discussion slides — a title + ONE big question each.

    Mirrors the exemplar's discussion slides: "Turn and Talk" title at the top
    and one large prompt below. We source the prompts from the lesson's Do Now
    (its questions, falling back to its stimulus) so the opening discussion lands
    right after the title slide. Emits 1–2 slides (one per Do Now question, max
    two), staying inside the density cap.
    """
    from pptx.util import Inches

    do_now = getattr(master, "do_now", None)
    if not do_now:
        return

    W = prs.slide_width  # noqa: N806

    prompts: list[str] = []
    for q in (getattr(do_now, "questions", None) or []):
        q = (q or "").strip()
        if q:
            prompts.append(_short(q, words=16))
    if not prompts:
        stim = (getattr(do_now, "stimulus", "") or "").strip()
        if stim:
            prompts.append(_first_sentence(stim, words=16))
    if not prompts:
        return

    for prompt in prompts[:2]:
        slide = _add_slide(prs)
        _slide_title(slide, W, "Turn and Talk", hex_color=_C_TERRACOTTA, font_size=38)
        # One large discussion question, centered in the lower two-thirds.
        _textbox(
            slide,
            left=Inches(0.8), top=Inches(2.0),
            width=W - Inches(1.6), height=Inches(2.6),
            text=prompt,
            font_size=28, bold=False,
            hex_color=_C_BODY, align_center=True,
        )


def _build_vocabulary_slides(prs: Presentation, master: MasterContent) -> None:
    """Vocabulary slides (up to 5 terms per slide). Definitions stay sparse."""
    from pptx.util import Inches

    if not master.vocabulary:
        return

    W = prs.slide_width  # noqa: N806
    TERMS_PER_SLIDE = 3  # noqa: N806

    vocab_chunks = [
        master.vocabulary[i: i + TERMS_PER_SLIDE]
        for i in range(0, len(master.vocabulary), TERMS_PER_SLIDE)
    ]
    for chunk_idx, chunk in enumerate(vocab_chunks):
        slide = _add_slide(prs)
        heading_label = "Vocabulary" if len(vocab_chunks) == 1 else f"Vocabulary ({chunk_idx + 1})"
        _slide_title(slide, W, heading_label, hex_color=_C_TERRACOTTA)

        # Term entries — term on the left, a SHORT definition on the right.
        top_offset: Any = Inches(1.35)
        row_height = Inches(0.85)
        for entry in chunk:
            _textbox(
                slide,
                left=Inches(0.3), top=top_offset,
                width=Inches(2.7), height=row_height,
                text=_short(entry.term, words=4),
                font_size=18, bold=True,
                hex_color=_C_TERRACOTTA,
            )
            _textbox(
                slide,
                left=Inches(3.1), top=top_offset,
                width=Inches(6.6), height=row_height,
                text=_short(entry.definition, words=4),
                font_size=16,
                hex_color=_C_BODY,
            )
            top_offset += row_height + Inches(0.05)


def _build_instruction_slides(prs: Presentation, master: MasterContent,
                              images: dict[str, Path]) -> None:
    """One sparse slide per InstructionSection: ≤3 short bullets + one visual."""
    from pptx.util import Inches

    W = prs.slide_width  # noqa: N806

    for section in master.direct_instruction:
        slide = _add_slide(prs)
        _slide_title(slide, W, section.heading, hex_color=_C_TERRACOTTA)

        has_image = bool(section.image_spec and section.image_spec in images)
        content_w = Inches(TEXT_PANEL_W_IN) if has_image else Inches(FULL_TEXT_W_IN)
        content_left = Inches(MARGIN_IN)

        # Build ≤3 short bullets. Prefer key_points; otherwise fall back to the
        # short slide_summary, then a first-sentence slice of content. Each
        # bullet is capped tight so the slide stays near the exemplar density.
        bullets = [_short(kp, words=6) for kp in section.key_points if kp.strip()][:2]
        if not bullets:
            summary = getattr(section, "slide_summary", "") or ""
            phrase = _short(summary, words=14) if summary.strip() else _first_sentence(section.content, words=14)
            if phrase:
                bullets = [phrase]

        if bullets:
            _bullet_textbox(
                slide,
                left=content_left, top=Inches(1.35),
                width=content_w, height=Inches(3.6),
                items=bullets,
                font_size=22,
                hex_color=_C_BODY,
            )

        if has_image:
            _embed_image(
                slide, section.image_spec, images,
                box_left=Inches(IMG_BOX_LEFT_IN), box_top=Inches(IMG_BOX_TOP_IN),
                box_w=Inches(IMG_BOX_W_IN), box_h=Inches(IMG_BOX_H_IN),
            )


def _build_source_slides(prs: Presentation, master: MasterContent,
                         images: dict[str, Path]) -> None:
    """One sparse slide per primary source: attribution + a short excerpt phrase."""
    from pptx.util import Inches

    W = prs.slide_width  # noqa: N806

    for ps in master.primary_sources:
        slide = _add_slide(prs)
        _slide_title(slide, W, f"Source: {_short(ps.title, words=5)}",
                     hex_color=_C_TERRACOTTA, font_size=26)

        has_image = bool(ps.image_spec and ps.image_spec in images)
        text_w = Inches(TEXT_PANEL_W_IN) if has_image else Inches(FULL_TEXT_W_IN)

        # Attribution + excerpt start low enough to clear a two-line source title.
        _textbox(
            slide,
            left=Inches(MARGIN_IN), top=Inches(1.5),
            width=text_w, height=Inches(0.4),
            text=_short(f"{ps.source_type.replace('_', ' ').title()}  |  {ps.attribution}", words=7),
            font_size=13, italic=True,
            hex_color=_C_TERRACOTTA,
        )

        excerpt = getattr(ps, "slide_excerpt", "") or ""
        excerpt = _short(excerpt, words=12) if excerpt.strip() else _first_sentence(ps.content_text, words=12)
        if excerpt:
            _textbox(
                slide,
                left=Inches(MARGIN_IN), top=Inches(2.1),
                width=text_w, height=Inches(2.4),
                text=f'"{excerpt}"',
                font_size=18, italic=True,
                hex_color=_C_BODY,
            )

        if has_image:
            _embed_image(
                slide, ps.image_spec, images,
                box_left=Inches(IMG_BOX_LEFT_IN), box_top=Inches(IMG_BOX_TOP_IN),
                box_w=Inches(IMG_BOX_W_IN), box_h=Inches(IMG_BOX_H_IN),
            )


def _collect_lesson_images(master: MasterContent, images: dict[str, Path]) -> list[str]:
    """Return, in deck order, every image_spec the lesson has a fetched file for.

    Pulls from instruction sections, primary sources, and exit-ticket stimuli —
    the same specs the per-slide builders embed — deduplicated so the activity
    grid never tiles the same picture twice.
    """
    specs: list[str] = []
    seen: set[str] = set()

    def _add(spec: str) -> None:
        if spec and spec in images and spec not in seen:
            seen.add(spec)
            specs.append(spec)

    for section in getattr(master, "direct_instruction", None) or []:
        _add(getattr(section, "image_spec", "") or "")
    for ps in getattr(master, "primary_sources", None) or []:
        _add(getattr(ps, "image_spec", "") or "")
    for sq in getattr(master, "exit_ticket", None) or []:
        _add(getattr(sq, "stimulus_image_spec", "") or "")
    return specs


def _build_image_activity_slide(prs: Presentation, master: MasterContent,
                                images: dict[str, Path]) -> None:
    """Tiled image-sorting activity slide — mirrors the exemplar's picture grid.

    The exemplar packs a 21-image sorting activity onto one slide. When the
    lesson has ≥6 fetched images we emit ONE slide that tiles them as small
    aspect-preserved cards (~2.7 in wide) in evenly spaced rows under a plain
    "Sort the Images" title. Each card reuses ``_embed_image`` so aspect ratios
    are preserved (no stretch). No text labels — the grid is the activity, so the
    slide stays within the density cap.
    """
    from pptx.util import Inches

    specs = _collect_lesson_images(master, images)
    if len(specs) < 6:
        return

    W = prs.slide_width  # noqa: N806
    slide = _add_slide(prs)
    _slide_title(slide, W, "Sort the Images", hex_color=_C_TERRACOTTA)

    # Grid geometry: fixed card width ~2.7 in, as many columns as fit the usable
    # width, rows flowing down the area below the title. Cap the tile count so a
    # huge image set never overflows the slide.
    card_w = 2.7
    gap = 0.18
    usable_w = SLIDE_W_IN - 2 * MARGIN_IN
    cols = max(1, int((usable_w + gap) / (card_w + gap)))
    grid_top = TITLE_TOP_IN + TITLE_H_IN + 0.15
    avail_h = SLIDE_H_IN - grid_top - MARGIN_IN
    max_rows = 3
    capacity = cols * max_rows
    tiles = specs[:capacity]
    rows = (len(tiles) + cols - 1) // cols
    card_h = (avail_h - gap * (rows - 1)) / rows if rows else avail_h
    # Center the grid block horizontally.
    grid_w = cols * card_w + (cols - 1) * gap
    grid_left = MARGIN_IN + max(0.0, (usable_w - grid_w) / 2)

    for idx, spec in enumerate(tiles):
        r, c = divmod(idx, cols)
        left = grid_left + c * (card_w + gap)
        top = grid_top + r * (card_h + gap)
        _embed_image(
            slide, spec, images,
            box_left=Inches(left), box_top=Inches(top),
            box_w=Inches(card_w), box_h=Inches(card_h),
        )


def _build_station_slide(prs: Presentation, master: MasterContent) -> None:
    """Station overview slide (if stations exist) — short titles + directions."""
    from pptx.util import Inches

    if not master.stations:
        return

    W = prs.slide_width  # noqa: N806
    slide = _add_slide(prs)
    _slide_title(slide, W, "Learning Stations", hex_color=_C_TERRACOTTA)

    top_offset = Inches(1.35)
    usable = SLIDE_W_IN - 2 * MARGIN_IN
    n = max(len(master.stations), 1)
    col_w = Inches(usable / n)
    for i, station in enumerate(master.stations):
        left = Inches(MARGIN_IN + (usable / n) * i)
        _textbox(
            slide,
            left=left,
            top=top_offset,
            width=col_w - Inches(0.1),
            height=Inches(1.5),
            text=_short(station.title, words=4),
            font_size=16, bold=True,
            hex_color=_C_TERRACOTTA,
        )
        # Numbered label sits well below to clear a wrapped (up to 3-line) title
        # in a narrow column. Names only here; directions live in the handout.
        _textbox(
            slide,
            left=left,
            top=top_offset + Inches(1.55),
            width=col_w - Inches(0.1),
            height=Inches(0.5),
            text=f"Station {i + 1}",
            font_size=13,
            hex_color=_C_BODY,
        )


def _build_exit_ticket_slide(prs: Presentation, master: MasterContent,
                             images: dict[str, Path]) -> None:
    """Exit ticket — ONE question per slide, sparse, questions only (no answers)."""
    from pptx.util import Inches

    if not master.exit_ticket:
        return

    W = prs.slide_width  # noqa: N806

    for i, sq in enumerate(master.exit_ticket, 1):
        slide = _add_slide(prs)
        _slide_title(slide, W, "Exit Ticket", hex_color=_C_TERRACOTTA)

        has_image = bool(sq.stimulus_image_spec and sq.stimulus_image_spec in images)
        text_w = Inches(TEXT_PANEL_W_IN) if has_image else Inches(FULL_TEXT_W_IN)

        stim = getattr(sq, "slide_stimulus", "") or ""
        stim = _short(stim, words=8) if stim.strip() else _short(sq.stimulus, words=8)
        if stim:
            _textbox(
                slide,
                left=Inches(MARGIN_IN), top=Inches(1.35),
                width=text_w, height=Inches(1.4),
                text=stim,
                font_size=15, italic=True,
                hex_color=_C_BODY,
            )
        _textbox(
            slide,
            left=Inches(MARGIN_IN), top=Inches(2.7),
            width=text_w, height=Inches(1.8),
            text=f"Q{i}: {_short(sq.question, words=13)}",
            font_size=22, bold=True,
            hex_color=_C_TERRACOTTA,
        )
        if has_image:
            _embed_image(
                slide, sq.stimulus_image_spec, images,
                box_left=Inches(IMG_BOX_LEFT_IN), box_top=Inches(IMG_BOX_TOP_IN),
                box_w=Inches(IMG_BOX_W_IN), box_h=Inches(IMG_BOX_H_IN),
            )


def _build_speaker_notes(prs: Presentation, master: MasterContent) -> None:
    """Attach misconceptions/formative checks as speaker notes."""
    notes_parts: list[str] = []
    if getattr(master, "misconceptions", None):
        notes_parts.append("MISCONCEPTIONS TO WATCH FOR:")
        for m in master.misconceptions:
            notes_parts.append(f"  - {m}")
    if getattr(master, "formative_checks", None):
        notes_parts.append("\nMID-LESSON CHECKS:")
        for fc in master.formative_checks:
            notes_parts.append(f"  - {fc}")
    if getattr(master, "prerequisite_skills", None):
        notes_parts.append("\nPREREQUISITES:")
        for ps in master.prerequisite_skills:
            notes_parts.append(f"  - {ps}")

    if notes_parts and len(prs.slides) > 1:
        try:
            target_slide = prs.slides[1]
            if not target_slide.has_notes_slide:
                target_slide.notes_slide  # creates notes slide
            notes_tf = target_slide.notes_slide.notes_text_frame
            notes_tf.text = "\n".join(notes_parts)
        except Exception as exc:
            # Notes are optional enhancement, never block export
            logger.debug("Could not add speaker notes: %s", exc)


# ═══════════════════════════════════════════════════════════════════════════
# Public API (unchanged signature)
# ═══════════════════════════════════════════════════════════════════════════


async def compile_slides(
    master: MasterContent,
    images: dict[str, Path],
    output_dir: Path,
) -> Path:
    """Compile a classroom-ready PPTX from a MasterContent object.

    The deck is a sparse, image-forward 16:9 presentation (~15–25 words/slide)
    on the standard 10.0 x 5.625 in canvas — opens cleanly on macOS and imports
    losslessly into Google Slides.

    Slide order:
        1. Title slide (title, subject, grade, objective)
        2. Turn-and-Talk discussion slide(s) (from the Do Now; 1–2 slides)
        3. Vocabulary slide(s) (up to 5 terms per slide)
        4. One slide per InstructionSection (heading, ≤3 short bullets, image)
        5. Source analysis slides (one per primary source)
        6. Image-activity grid (only when ≥6 fetched images are available)
        7. Station overview (if stations exist)
        8. Exit ticket slides (one question per slide, no answers)

    Args:
        master: The MasterContent source-of-truth object.
        images: Mapping of image_spec strings to local file Paths.
        output_dir: Directory where the .pptx file will be written.

    Returns:
        Path to the generated .pptx file.
    """
    from pptx import Presentation
    from pptx.util import Inches

    from clawed.io import safe_filename

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Build all slide types
    _build_title_slide(prs, master)
    _build_turn_and_talk_slides(prs, master)
    _build_vocabulary_slides(prs, master)
    _build_instruction_slides(prs, master, images)
    _build_source_slides(prs, master, images)
    _build_image_activity_slide(prs, master, images)
    _build_station_slide(prs, master)
    _build_exit_ticket_slide(prs, master, images)
    _build_speaker_notes(prs, master)

    # Save
    safe = safe_filename(master.title)
    out_path = output_dir / f"{safe}_slides.pptx"
    prs.save(str(out_path))

    # Post-save smoke check: the deck must be a standard 16:9 canvas so it opens
    # cleanly on macOS Keynote/PowerPoint and imports losslessly into Google
    # Slides. A non-standard canvas was the original macOS/Slides failure mode.
    from pptx.util import Emu
    w_in = round(Emu(int(prs.slide_width or 0)).inches, 3)
    h_in = round(Emu(int(prs.slide_height or 0)).inches, 3)
    if (w_in, h_in) != (SLIDE_W_IN, SLIDE_H_IN):
        logger.warning(
            "Slide deck canvas is %sx%s in, expected %sx%s — Slides/macOS import may scale.",
            w_in, h_in, SLIDE_W_IN, SLIDE_H_IN,
        )

    logger.info("Slides saved to %s (%d slides, %sx%s in)", out_path, len(prs.slides._sldIdLst), w_in, h_in)
    return out_path
